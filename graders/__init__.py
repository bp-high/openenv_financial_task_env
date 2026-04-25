"""Grading functions for the Financial Task Environment.

Two grading modes:
  1. QA tasks  — compare agent text answer against reference text
                 (numeric extraction + keyword matching)
  2. MODIFY tasks — compare agent-produced xlsx against reference xlsx
                    (cell-level comparison with tolerance)
"""

from __future__ import annotations

import re
import traceback
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence

import openpyxl
from openpyxl.utils import get_column_letter


# ---------------------------------------------------------------------------
# Numeric helpers
# ---------------------------------------------------------------------------

def _extract_numbers(text: str) -> List[float]:
    """Extract all numeric values from text, handling commas, $, %."""
    cleaned = text.replace("$", "").replace("€", "").replace("£", "")
    pattern = r"-?\d{1,3}(?:,\d{3})*(?:\.\d+)?|-?\d+(?:\.\d+)?"
    raw = re.findall(pattern, cleaned)
    results: List[float] = []
    for r in raw:
        try:
            results.append(float(r.replace(",", "")))
        except ValueError:
            continue
    return results


def _number_close(actual: float, expected: float, rel_tol: float = 0.05) -> bool:
    if expected == 0:
        return abs(actual) < 1e-6
    return abs(actual - expected) / abs(expected) <= rel_tol


def _best_number_match(numbers: List[float], target: float, rel_tol: float = 0.05) -> bool:
    return any(_number_close(n, target, rel_tol) for n in numbers)


# ---------------------------------------------------------------------------
# QA grading (text answer)
# ---------------------------------------------------------------------------

def grade_qa(answer: str, reference_answer: str) -> float:
    """Grade a text answer against a reference.  Returns 0.0–1.0."""
    if not answer.strip():
        return 0.0

    ref_nums = _extract_numbers(reference_answer)
    ans_nums = _extract_numbers(answer)

    if ref_nums:
        # Numeric comparison: what fraction of reference numbers appear?
        matched = sum(1 for r in ref_nums if _best_number_match(ans_nums, r))
        num_score = matched / len(ref_nums)
    else:
        num_score = 0.0

    # Keyword overlap
    ref_words = set(re.findall(r"[a-zA-Z]{3,}", reference_answer.lower()))
    ans_words = set(re.findall(r"[a-zA-Z]{3,}", answer.lower()))
    if ref_words:
        kw_score = len(ref_words & ans_words) / len(ref_words)
    else:
        kw_score = 0.0

    # Weighted combination (numbers matter more for financial tasks)
    if ref_nums:
        # If all numbers match perfectly, give full score
        if num_score >= 1.0:
            return 1.0
        return round(min(0.8 * num_score + 0.2 * kw_score, 1.0), 4)
    else:
        return round(kw_score, 4)


# ---------------------------------------------------------------------------
# MODIFY grading (xlsx comparison)
# ---------------------------------------------------------------------------

def _load_wb_values(path: str):
    """Load workbook in data_only mode, return dict of {(sheet, row, col): value}."""
    wb = openpyxl.load_workbook(path, data_only=True)
    cells = {}
    sheets = set()
    for name in wb.sheetnames:
        sheets.add(name)
        ws = wb[name]
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None:
                    cells[(name, cell.row, cell.column)] = cell.value
    wb.close()
    return cells, sheets


def grade_xlsx(output_path: str, reference_path: str) -> float:
    """Compare agent output xlsx with reference xlsx.  Returns 0.0–1.0.

    Scoring breakdown:
      - 30%  sheet-level: does the output have all reference sheets?
      - 70%  cell-level:  fraction of reference cells matched (with tolerance for numbers)
    """
    try:
        ref_cells, ref_sheets = _load_wb_values(reference_path)
        out_cells, out_sheets = _load_wb_values(output_path)
    except Exception:
        return 0.0

    # --- Sheet score (30%) ---
    if ref_sheets:
        sheet_score = len(ref_sheets & out_sheets) / len(ref_sheets)
    else:
        sheet_score = 1.0

    # --- Cell score (70%) ---
    if not ref_cells:
        return round(0.3 * sheet_score + 0.7 * 1.0, 4)

    matched = 0
    total = len(ref_cells)

    for key, ref_val in ref_cells.items():
        out_val = out_cells.get(key)
        if out_val is None:
            continue
        if ref_val == out_val:
            matched += 1
            continue
        # Numeric tolerance
        try:
            rv = float(ref_val)
            ov = float(out_val)
            if _number_close(ov, rv, rel_tol=0.02):
                matched += 1
                continue
        except (ValueError, TypeError):
            pass
        # String comparison (case-insensitive, whitespace-normalized)
        try:
            if str(ref_val).strip().lower() == str(out_val).strip().lower():
                matched += 1
        except Exception:
            pass

    cell_score = matched / total if total > 0 else 1.0

    return round(0.3 * sheet_score + 0.7 * cell_score, 4)


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

def _clamp_score(score: float) -> float:
    """Clamp score to open interval (0.001, 0.999) — evaluator rejects exact 0.0 and 1.0."""
    return max(0.001, min(0.999, score))


# ---------------------------------------------------------------------------
# DOCX grading — 3-layer composition
#   1. Validity gate    file parses with python-docx           (binary, gates the rest)
#   2. Generic diff     paragraph-level fuzzy compare to gold  (0–1)
#   3. Per-task         OSWorld evaluator function(s)          (0–1)
# Final = clamp( 0.4 * diff + 0.6 * primary )  iff validity gate passes.
# ---------------------------------------------------------------------------

def _docx_validity(path: str) -> bool:
    try:
        from docx import Document
        Document(path)
        return True
    except Exception:
        return False


def _docx_diff(working: str, gold: str) -> float:
    """Generic paragraph-level fuzzy ratio against the primary gold."""
    if not gold:
        return 0.0
    try:
        from graders.docx_metrics import compare_docx_files
        return float(compare_docx_files(working, gold, ignore_blanks=True, fuzzy_match=True))
    except Exception:
        return 0.0


def grade_docx(task: Dict[str, Any], output_path: str) -> float:
    if not _docx_validity(output_path):
        return 0.001  # corruption / not a docx → fail outright

    # Layer 2: generic diff vs the first gold file
    primary_gold = task.get("reference_file", "")
    diff_score = _docx_diff(output_path, primary_gold) if primary_gold else 0.0

    # Layer 3: per-task evaluator (OSWorld checks)
    evaluator = task.get("evaluator") or {}
    conj = evaluator.get("conj", "and")
    checks = evaluator.get("checks") or []
    if checks:
        from graders.docx_metrics import run_evaluator
        # Resolve relative paths in expected_files against the repo root.
        # __file__ is graders/__init__.py → parent is graders/, parent.parent is repo root.
        from pathlib import Path as _P
        repo_root = _P(__file__).resolve().parent.parent
        resolved_checks: List[Dict[str, Any]] = []
        for c in checks:
            ef = c.get("expected_files") or []
            ef_abs = [str(repo_root / f) if not _P(f).is_absolute() else f for f in ef]
            resolved_checks.append({**c, "expected_files": ef_abs})
        primary_score = run_evaluator(
            conj=conj,
            checks=resolved_checks,
            working_file=output_path,
            source_file=task.get("source_file", ""),
        )
    else:
        primary_score = diff_score  # no per-task evaluator → diff carries it all

    return round(0.4 * diff_score + 0.6 * primary_score, 4)


# ---------------------------------------------------------------------------
# PPTX grading — 2-layer composition
#   1. Validity gate    file parses with python-pptx           (binary, gates the rest)
#   2. Structural diff  slide-count match (30%)
#                       + per-(slide_idx, shape_idx) text match (70%)
# Per-task evaluator (PPTArena's VLM judge) is intentionally NOT wired:
# it's expensive and non-deterministic.  See edits.md for plan.
# ---------------------------------------------------------------------------

def _pptx_validity(path: str) -> bool:
    try:
        from pptx import Presentation
        Presentation(path)
        return True
    except Exception:
        return False


def _shape_style(shape) -> Dict[str, Any]:
    """Extract per-shape style attributes.  All values are tolerant — any
    attribute that can't be read (placeholder inheritance, unsupported shape
    type, None color, exception) becomes None.  Two None values on the same
    key match; one None vs one non-None counts as a mismatch."""
    style: Dict[str, Any] = {
        "fill_rgb":       None,  # solid fill RGB hex (str)
        "fill_theme":     None,  # theme color name (str) — captures gold-uses-theme tasks
        "line_rgb":       None,  # line/border RGB hex (str)
        "font_name":      None,  # first run, str
        "font_size_pt":   None,  # first run, float
        "font_bold":      None,  # first run, bool/None (None = inherited)
        "font_italic":    None,  # first run, bool/None
        "font_rgb":       None,  # first run, RGB hex (str)
        "para_alignment": None,  # first paragraph: PP_ALIGN.CENTER, LEFT, RIGHT, JUSTIFY (str)
    }
    # ---- shape fill / line colors ----
    # Capture both explicit RGB and theme-color reference: a "match colors to
    # theme" task changes ad-hoc RGB to theme references, and we want to see
    # that as a real difference rather than letting both sides become None.
    try:
        from pptx.enum.dml import MSO_FILL_TYPE
        ft = shape.fill.type
        if ft == MSO_FILL_TYPE.SOLID:
            try:
                style["fill_rgb"] = str(shape.fill.fore_color.rgb)
            except Exception:
                # solid fill but rgb raises → it's a theme color
                try:
                    style["fill_theme"] = str(shape.fill.fore_color.theme_color)
                except Exception:
                    pass
    except Exception:
        pass
    try:
        style["line_rgb"] = str(shape.line.color.rgb)
    except Exception:
        pass
    # ---- first-paragraph + first-run text properties ----
    try:
        if shape.has_text_frame:
            tf = shape.text_frame
            if tf.paragraphs:
                p0 = tf.paragraphs[0]
                # Paragraph-level alignment (LEFT/CENTER/RIGHT/JUSTIFY) —
                # critical for "center the title", "right-align" tasks.
                try:
                    style["para_alignment"] = str(p0.alignment) if p0.alignment is not None else None
                except Exception:
                    pass
                if p0.runs:
                    run = p0.runs[0]
                    style["font_name"] = run.font.name
                    if run.font.size is not None:
                        style["font_size_pt"] = float(run.font.size.pt)
                    style["font_bold"] = run.font.bold
                    style["font_italic"] = run.font.italic
                    try:
                        style["font_rgb"] = str(run.font.color.rgb)
                    except Exception:
                        pass
    except Exception:
        pass
    return style


def _pptx_load_slides(path: str):
    """Returns (slide_w, slide_h, [ [shape_dict, ...] per slide ])
    where shape_dict has: text, left, top, width, height (None if inherited),
    and `style` (per-attribute dict from _shape_style).
    Dimensions are in EMU (English Metric Units, python-pptx native)."""
    from pptx import Presentation
    prs = Presentation(path)
    slide_w = prs.slide_width or 9144000   # default 10in if unset
    slide_h = prs.slide_height or 6858000  # default 7.5in
    slides = []
    for slide in prs.slides:
        shapes = []
        for shape in slide.shapes:
            try:
                text = shape.text_frame.text.strip() if shape.has_text_frame else ""
            except Exception:
                text = ""
            shapes.append({
                "text": text,
                "left": getattr(shape, "left", None),
                "top": getattr(shape, "top", None),
                "width": getattr(shape, "width", None),
                "height": getattr(shape, "height", None),
                "style": _shape_style(shape),
            })
        slides.append(shapes)
    return slide_w, slide_h, slides


def _coord_match(r_val, o_val, denom: int) -> float:
    """Tolerance-based score for one coordinate (left/top/width/height).

    Both None  → 1.0  (placeholder inheriting from layout, treated as match)
    Either None or denom 0 → 0.0
    delta <= 2% of denom  → 1.0  (visually indistinguishable)
    delta >= 20% of denom → 0.0  (clearly different)
    in between → linear falloff
    """
    if r_val is None and o_val is None:
        return 1.0
    if r_val is None or o_val is None or not denom:
        return 0.0
    delta = abs(r_val - o_val) / denom
    if delta <= 0.02:
        return 1.0
    if delta >= 0.20:
        return 0.0
    return 1.0 - (delta - 0.02) / 0.18  # linear 1.0 → 0.0 over 2%–20%


_STYLE_WEIGHTS = {
    "fill_rgb":       0.22,  # most often-edited attribute in styling tasks
    "fill_theme":     0.08,  # NEW: theme-color reference (catches "match colors to theme")
    "line_rgb":       0.08,
    "font_name":      0.08,
    "font_size_pt":   0.12,
    "font_bold":      0.05,
    "font_italic":    0.05,
    "font_rgb":       0.17,
    "para_alignment": 0.15,  # NEW: catches "center the title" and friends
}


def _style_match_score(ref_style: Dict[str, Any], out_style: Dict[str, Any]) -> float:
    """Weighted attribute-by-attribute match.  Both None on a key counts as
    match (both inheriting from layout/theme).  Partial credit: matches sum,
    mismatches don't — score = sum(weighted matches) / sum(weights).
    """
    total_w = 0.0
    score = 0.0
    for key, w in _STYLE_WEIGHTS.items():
        r = ref_style.get(key)
        o = out_style.get(key)
        total_w += w
        if r == o:  # works for both None==None and value equality
            score += w
    return (score / total_w) if total_w > 0 else 1.0


def _shape_match_score(ref: dict, out: dict, slide_w: int, slide_h: int) -> float:
    """Per-shape composite: 40% text + 20% style + 20% position + 20% size."""
    # Text: exact match → 1.0; rapidfuzz fallback (ratio/100) so partial edits
    # get partial credit instead of binary fail.
    if ref["text"] == out["text"]:
        text_score = 1.0
    elif ref["text"] and out["text"]:
        try:
            from rapidfuzz import fuzz
            text_score = fuzz.ratio(ref["text"], out["text"]) / 100.0
        except Exception:
            text_score = 0.0
    else:
        text_score = 0.0

    style_score = _style_match_score(ref.get("style") or {}, out.get("style") or {})

    pos_score = (
        _coord_match(ref["left"], out["left"], slide_w)
        + _coord_match(ref["top"], out["top"], slide_h)
    ) / 2.0
    size_score = (
        _coord_match(ref["width"], out["width"], slide_w)
        + _coord_match(ref["height"], out["height"], slide_h)
    ) / 2.0

    return (
        0.40 * text_score
        + 0.20 * style_score
        + 0.20 * pos_score
        + 0.20 * size_score
    )


def grade_pptx(task: Dict[str, Any], output_path: str) -> float:
    if not _pptx_validity(output_path):
        return 0.001

    ref_path = task.get("reference_file", "")
    if not ref_path or not Path(ref_path).exists():
        return 0.001

    try:
        out_w, out_h, out_slides = _pptx_load_slides(output_path)
        ref_w, ref_h, ref_slides = _pptx_load_slides(ref_path)
    except Exception:
        return 0.001

    # Slide-count score (20%)
    if not ref_slides:
        slide_score = 1.0
    else:
        delta = abs(len(out_slides) - len(ref_slides))
        slide_score = max(0.0, 1.0 - (delta / max(len(ref_slides), 1)))

    # Per-shape avg score (80%) — text + position + size, weighted internally
    shape_scores: List[float] = []
    # Use ref's slide dimensions for normalization (the gold's coordinate frame)
    for s_i, ref_shapes in enumerate(ref_slides):
        out_shapes = out_slides[s_i] if s_i < len(out_slides) else []
        for sh_i, ref_shape in enumerate(ref_shapes):
            if sh_i < len(out_shapes):
                shape_scores.append(
                    _shape_match_score(ref_shape, out_shapes[sh_i], ref_w, ref_h)
                )
            else:
                shape_scores.append(0.0)  # missing shape in agent's output
    avg_shape = sum(shape_scores) / len(shape_scores) if shape_scores else 1.0

    return round(0.2 * slide_score + 0.8 * avg_shape, 4)


def _same_bytes(a: str, b: str) -> bool:
    """True iff two files exist and have identical SHA-256.  Used to detect
    'submitted source unchanged' exploits — a model that bypasses the work
    by handing back the input file."""
    import hashlib
    try:
        ah = hashlib.sha256(open(a, "rb").read()).hexdigest()
        bh = hashlib.sha256(open(b, "rb").read()).hexdigest()
        return ah == bh
    except Exception:
        return False


def grade_task(task: Dict[str, Any], answer: str = "", output_path: str = "") -> float:
    """Grade a task.  Returns score in (0.001, 0.999).

    For QA tasks:    uses *answer* (text) vs task["reference_answer"].
    For MODIFY tasks (xlsx): cell-diff against task["reference_file"].
    For MODIFY tasks (docx): validity + diff + OSWorld evaluator.
    For MODIFY tasks (pptx): validity + slide-count + per-shape composite.
    """
    task_type = task.get("task_type", "QA")
    family = task.get("family", "xlsx")

    if task_type == "QA":
        ref = task.get("reference_answer", "")
        return _clamp_score(grade_qa(answer, ref))
    elif task_type == "MODIFY":
        if not output_path or not Path(output_path).exists():
            return 0.001

        # ANTI-EXPLOIT: detect "submitted source unchanged".  A model that
        # discovers a task where source-vs-gold scores high (e.g., a small
        # alignment edit on a 55-shape deck) can game the diff by handing
        # back the input file.  We refuse to credit byte-identical
        # submissions UNLESS the task is OSWorld's `infeasible` sentinel
        # (where not-modifying is the correct answer).
        src = task.get("source_file", "")
        is_infeasible = False
        ev = task.get("evaluator") or {}
        for c in ev.get("checks") or []:
            if c.get("func") == "infeasible":
                is_infeasible = True
                break
        if src and Path(src).exists() and not is_infeasible:
            if _same_bytes(output_path, src):
                return 0.001

        if family == "docx":
            return _clamp_score(grade_docx(task, output_path))
        if family == "pptx":
            return _clamp_score(grade_pptx(task, output_path))
        # xlsx (default) and any future families fall back to cell-diff
        ref_path = task.get("reference_file", "")
        if not ref_path or not Path(ref_path).exists():
            return 0.001
        return _clamp_score(grade_xlsx(output_path, ref_path))
    else:
        return 0.001
