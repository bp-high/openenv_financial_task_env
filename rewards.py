"""Unified reward shaper for the office-document task environment.

Produces dense per-step process rewards across the xlsx / pptx / docx families.
Rewards are measured from real file-state changes, not from code-string heuristics.

Usage:
    tracker = RewardTracker(
        family="docx",
        working_file=path,
        gold_file=ref,
        enable_progress=True,
        task_evaluator=run_eval_fn,   # optional: callable[str -> float in 0–1]
    )
    # after each code step:
    signals = tracker.score_step(code=code, succeeded=ok, stdout=out)
    reward = signals.total  # bounded to [0, 0.10]

Components (all per-step, summed and clamped to 0.10):
    exec_health     code ran cleanly and produced output            [0.000–0.020]
    lib_engagement  code uses the family's expected library         [0.000–0.010]
    mutation        working file's hash changed this step           [0.000–0.030]
    validity        mutated file still parses for the family        [0.000–0.020]
    progress        structural distance to gold decreased           [0.000–0.040]
    eval_check      per-task evaluator score went UP this step      [0.000–0.020]

`progress` is gated: requires a gold file AND `enable_progress=True`.  Disable
for eval to keep the signal honest; enable during training for dense gradient.

`eval_check` requires `task_evaluator` to be passed.  It computes the per-task
evaluator (e.g. OSWorld docx checks) before/after a mutating step and rewards
*increases* — so the agent gets feedback when a previously-failing property
check starts passing.  Hardens against generic-distance-only gaming.
"""

from __future__ import annotations

import hashlib
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Optional

# ---------------------------------------------------------------------------
# Constants — keep total cap == 0.10 to stay backward-compatible with cumulative bounds
# ---------------------------------------------------------------------------

EXEC_HEALTH_FAIL = 0.005
EXEC_HEALTH_OK = 0.005
EXEC_HEALTH_OK_WITH_OUTPUT = 0.020
LIB_ENGAGEMENT = 0.010
MUTATION = 0.030
VALIDITY = 0.020
PROGRESS_MAX = 0.040
EVAL_CHECK_MAX = 0.020
STEP_CAP = 0.10

# Per-family library-detection regexes
_LIB_PATTERNS = {
    "xlsx": re.compile(r"\bopenpyxl\b|\bload_workbook\b|\bWorkbook\b"),
    "pptx": re.compile(r"\bpython-pptx\b|\bfrom\s+pptx\b|\bimport\s+pptx\b|\bPresentation\b"),
    "docx": re.compile(r"\bpython-docx\b|\bfrom\s+docx\b|\bimport\s+docx\b|\bDocument\b"),
}


# ---------------------------------------------------------------------------
# Step signals
# ---------------------------------------------------------------------------

@dataclass
class StepSignals:
    """Per-step reward decomposition — useful for both training and debugging."""

    exec_health: float = 0.0
    lib_engagement: float = 0.0
    mutation: float = 0.0
    validity: float = 0.0
    progress: float = 0.0
    eval_check: float = 0.0

    @property
    def total(self) -> float:
        s = (self.exec_health + self.lib_engagement + self.mutation
             + self.validity + self.progress + self.eval_check)
        return round(min(STEP_CAP, s), 4)

    def to_dict(self) -> dict:
        return {
            "exec_health": self.exec_health,
            "lib_engagement": self.lib_engagement,
            "mutation": self.mutation,
            "validity": self.validity,
            "progress": self.progress,
            "eval_check": self.eval_check,
            "total": self.total,
        }


# ---------------------------------------------------------------------------
# Tracker — one instance per episode
# ---------------------------------------------------------------------------

class RewardTracker:
    def __init__(
        self,
        family: str,
        working_file: str,
        gold_file: Optional[str] = None,
        enable_progress: bool = False,
        task_evaluator: Optional[Callable[[str], float]] = None,
    ) -> None:
        if family not in _LIB_PATTERNS:
            raise ValueError(f"unknown family {family!r}; expected xlsx|pptx|docx")
        self.family = family
        self.working_file = Path(working_file)
        self.gold_file = Path(gold_file) if gold_file else None
        self.enable_progress = bool(enable_progress and self.gold_file and self.gold_file.exists())
        self.task_evaluator = task_evaluator

        self._prev_hash = _file_hash(self.working_file)
        self._prev_distance: Optional[float] = (
            _structural_distance(self.family, self.working_file, self.gold_file)
            if self.enable_progress
            else None
        )
        # Baseline per-task evaluator score (e.g. OSWorld checks at episode start).
        # We only reward *increases* over this baseline.
        self._prev_eval: Optional[float] = (
            self._safe_task_eval() if self.task_evaluator is not None else None
        )

    def _safe_task_eval(self) -> float:
        try:
            return float(self.task_evaluator(str(self.working_file)))  # type: ignore[misc]
        except Exception:
            return 0.0

    # --------------------------------------------------------------
    def score_step(self, *, code: str, succeeded: bool, stdout: str) -> StepSignals:
        sig = StepSignals()

        # 1. Exec health — failed code gets minimal, success scales with output
        if not succeeded:
            sig.exec_health = EXEC_HEALTH_FAIL
            # Other signals stay 0; failure short-circuits.
            return sig
        sig.exec_health = EXEC_HEALTH_OK_WITH_OUTPUT if stdout.strip() else EXEC_HEALTH_OK

        # 2. Library engagement
        if _LIB_PATTERNS[self.family].search(code):
            sig.lib_engagement = LIB_ENGAGEMENT

        # 3. Mutation — file hash changed
        cur_hash = _file_hash(self.working_file)
        mutated = cur_hash != self._prev_hash
        if mutated:
            sig.mutation = MUTATION

            # 4. Validity — only meaningful if the file changed
            file_valid = _is_valid(self.family, self.working_file)
            if file_valid:
                sig.validity = VALIDITY

            # 5. Progress — structural-distance-to-gold decreased
            if self.enable_progress:
                cur_dist = _structural_distance(self.family, self.working_file, self.gold_file)
                if cur_dist is not None and self._prev_distance is not None:
                    delta = self._prev_distance - cur_dist
                    if delta > 0:
                        # Scale: up to PROGRESS_MAX as agent closes the gap.
                        # Normalize by the *initial* distance so big leaps near the
                        # start don't dominate small refinements near the end.
                        denom = max(self._prev_distance, 0.05)
                        sig.progress = min(PROGRESS_MAX, PROGRESS_MAX * (delta / denom))
                if cur_dist is not None:
                    self._prev_distance = cur_dist

            # 6. Per-task evaluator — reward increases in the spec-aligned score.
            # Only fires if the file is still valid (running OSWorld evaluators
            # on a corrupt docx is wasted work).
            if self.task_evaluator is not None and file_valid:
                cur_eval = self._safe_task_eval()
                if self._prev_eval is not None and cur_eval > self._prev_eval:
                    delta = cur_eval - self._prev_eval
                    sig.eval_check = min(EVAL_CHECK_MAX, EVAL_CHECK_MAX * delta)
                self._prev_eval = cur_eval

        self._prev_hash = cur_hash
        return sig


# ---------------------------------------------------------------------------
# File hashing
# ---------------------------------------------------------------------------

def _file_hash(path: Path) -> str:
    if not path.exists():
        return ""
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


# ---------------------------------------------------------------------------
# Validity checks (per family)
# ---------------------------------------------------------------------------

def _is_valid(family: str, path: Path) -> bool:
    try:
        if family == "xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            wb.close()
            return True
        if family == "pptx":
            from pptx import Presentation
            Presentation(str(path))
            return True
        if family == "docx":
            from docx import Document
            Document(str(path))
            return True
    except Exception:
        return False
    return False


# ---------------------------------------------------------------------------
# Structural distance (per family) — 0.0 = identical, 1.0 = unrelated
# ---------------------------------------------------------------------------

def _structural_distance(family: str, working: Path, gold: Path) -> Optional[float]:
    try:
        if family == "xlsx":
            return _xlsx_distance(working, gold)
        if family == "pptx":
            return _pptx_distance(working, gold)
        if family == "docx":
            return _docx_distance(working, gold)
    except Exception:
        return None
    return None


def _xlsx_distance(a: Path, b: Path) -> float:
    """1 - (fraction of gold cells matched in working).  Cheap, mirrors grader."""
    import openpyxl

    def _load(p: Path) -> dict:
        wb = openpyxl.load_workbook(p, data_only=True, read_only=True)
        cells: dict = {}
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=False):
                for c in row:
                    if c.value is not None:
                        cells[(sheet, c.row, c.column)] = c.value
        wb.close()
        return cells

    a_cells = _load(a)
    b_cells = _load(b)
    if not b_cells:
        return 1.0
    matched = 0
    for k, v in b_cells.items():
        av = a_cells.get(k)
        if av is None:
            continue
        if av == v:
            matched += 1
            continue
        try:
            if abs(float(av) - float(v)) / max(abs(float(v)), 1e-9) <= 0.02:
                matched += 1
                continue
        except (TypeError, ValueError):
            pass
        if str(av).strip().lower() == str(v).strip().lower():
            matched += 1
    return 1.0 - (matched / len(b_cells))


def _pptx_distance(a: Path, b: Path) -> float:
    """1 - (fraction of gold shapes' text matched on same (slide, idx))."""
    from pptx import Presentation

    def _load(p: Path) -> dict:
        prs = Presentation(str(p))
        out: dict = {}
        for s_i, slide in enumerate(prs.slides):
            for sh_i, shape in enumerate(slide.shapes):
                txt = getattr(shape, "text_frame", None)
                txt = shape.text_frame.text if txt is not None else ""
                out[(s_i, sh_i)] = txt
        return out

    a_shapes = _load(a)
    b_shapes = _load(b)
    if not b_shapes:
        return 1.0
    matched = sum(1 for k, v in b_shapes.items() if a_shapes.get(k, "").strip() == v.strip())
    return 1.0 - (matched / len(b_shapes))


def _docx_distance(a: Path, b: Path) -> float:
    """1 - (fraction of gold paragraphs matched at same index)."""
    from docx import Document

    def _paras(p: Path) -> list[str]:
        doc = Document(str(p))
        return [para.text.strip() for para in doc.paragraphs]

    a_p = _paras(a)
    b_p = _paras(b)
    if not b_p:
        return 1.0
    matched = sum(1 for i, t in enumerate(b_p) if i < len(a_p) and a_p[i] == t)
    return 1.0 - (matched / len(b_p))
